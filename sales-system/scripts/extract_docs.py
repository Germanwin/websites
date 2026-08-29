import os, sys, json, pathlib
SRC = pathlib.Path(sys.argv[1] if len(sys.argv) > 1 else "/Users/hermmoment/coding/websites/sales-system/00-source")
OUT = pathlib.Path("/Users/hermmoment/coding/websites/sales-system/02-docs-text")
OUT.mkdir(parents=True, exist_ok=True)

def flat(p):
    return str(p.relative_to(SRC)).replace("/", "__")

report = []

def pdf(p):
    from pypdf import PdfReader
    r = PdfReader(str(p))
    out=[]
    for i,pg in enumerate(r.pages,1):
        t = (pg.extract_text() or "").strip()
        out.append(f"--- стр. {i} ---\n{t}")
    return "\n\n".join(out), len(r.pages)

def pptx(p):
    from pptx import Presentation
    pr = Presentation(str(p)); out=[]
    for i,s in enumerate(pr.slides,1):
        parts=[]
        for sh in s.shapes:
            if sh.has_text_frame:
                tx = sh.text_frame.text.strip()
                if tx: parts.append(tx)
            if sh.has_table:
                for row in sh.table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
        out.append(f"--- слайд {i} ---\n" + "\n".join(parts))
    return "\n\n".join(out), len(pr.slides)

def xlsx(p):
    import openpyxl
    wb = openpyxl.load_workbook(str(p), data_only=True); out=[]
    for ws in wb.worksheets:
        out.append(f"--- лист: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells=[str(c) for c in row if c is not None]
            if cells: out.append(" | ".join(cells))
    return "\n".join(out), len(wb.worksheets)

def docx(p):
    import docx as D
    d = D.Document(str(p)); out=[para.text for para in d.paragraphs if para.text.strip()]
    for t in d.tables:
        for row in t.rows:
            out.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(out), len(out)

H = {".pdf":pdf, ".pptx":pptx, ".xlsx":xlsx, ".docx":docx}

for p in sorted(SRC.rglob("*")):
    if not p.is_file(): continue
    ext = p.suffix.lower()
    if ext == ".txt":
        txt = p.read_text(encoding="utf-8", errors="replace")
        (OUT / (flat(p) + ".md")).write_text(f"# {p.relative_to(SRC)}\n\n{txt}", encoding="utf-8")
        report.append((str(p.relative_to(SRC)), "txt", len(txt), "ok"))
        continue
    if ext not in H: continue
    try:
        txt, n = H[ext](p)
        (OUT / (flat(p) + ".md")).write_text(f"# {p.relative_to(SRC)}\n\n{txt}", encoding="utf-8")
        status = "ok" if len(txt.replace("---","").strip()) > 100 else "EMPTY/SCANNED"
        report.append((str(p.relative_to(SRC)), ext, len(txt), status))
    except Exception as e:
        report.append((str(p.relative_to(SRC)), ext, 0, f"ERROR {e}"))

print(f"{'STATUS':16} {'CHARS':>7}  FILE")
for f,e,n,s in report:
    print(f"{s:16} {n:>7}  {f}")
print("\nTOTAL:", len(report))
